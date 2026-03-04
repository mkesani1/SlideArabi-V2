#!/usr/bin/env python3
"""
test_harness.py — SlideShift v2 Structural Comparison Test Harness

Runs structural analysis and transformation on one or more input PPTX files,
then generates per-deck HTML reports comparing LTR→RTL structure.

Usage:
    python test_harness.py input1.pptx [input2.pptx ...]
    python test_harness.py --dir /path/to/decks/
    python test_harness.py input.pptx --output-dir /path/to/reports/

Outputs (per deck):
    {deck_name}_report.html   — side-by-side structural diff
    {deck_name}_rtl.pptx      — transformed RTL output
    run_summary.json          — aggregate stats across all decks
"""
from __future__ import annotations

import argparse
import json
import logging
import os
import sys
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from slideshift_v2.layout_analyzer import LayoutAnalyzer
    from slideshift_v2.structural_validator import StructuralValidator
    from slideshift_v2.rtl_transforms import MasterLayoutTransformer, SlideContentTransformer
    from slideshift_v2.property_resolver import PropertyResolver
    from slideshift_v2.template_registry import TemplateRegistry
except ImportError as e:
    print(f"WARNING: Some v2 modules not available: {e}")
    print("Continuing with available modules...")

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Data Classes
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class ShapeSnapshot:
    """Snapshot of a single shape's structural properties."""
    name: str
    shape_type: str
    left: int
    top: int
    width: int
    height: int
    has_text: bool
    text_preview: str = ""
    is_placeholder: bool = False
    placeholder_type: Optional[str] = None
    placeholder_idx: Optional[int] = None
    rtl_flag: Optional[bool] = None
    alignment: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "name": self.name,
            "type": self.shape_type,
            "left": self.left,
            "top": self.top,
            "width": self.width,
            "height": self.height,
            "has_text": self.has_text,
            "text_preview": self.text_preview,
            "is_placeholder": self.is_placeholder,
            "placeholder_type": self.placeholder_type,
            "placeholder_idx": self.placeholder_idx,
            "rtl_flag": self.rtl_flag,
            "alignment": self.alignment,
        }


@dataclass
class SlideSnapshot:
    """Snapshot of a single slide's structural state."""
    slide_num: int
    layout_name: str
    layout_type: str
    shape_count: int
    shapes: List[ShapeSnapshot] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "slide_num": self.slide_num,
            "layout_name": self.layout_name,
            "layout_type": self.layout_type,
            "shape_count": self.shape_count,
            "shapes": [s.to_dict() for s in self.shapes],
        }


@dataclass
class DeckSnapshot:
    """Snapshot of an entire presentation's structural state."""
    file_path: str
    slide_count: int
    slide_width: int
    slide_height: int
    slides: List[SlideSnapshot] = field(default_factory=list)
    capture_timestamp: float = field(default_factory=time.time)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_path": self.file_path,
            "slide_count": self.slide_count,
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
            "slides": [s.to_dict() for s in self.slides],
            "capture_timestamp": self.capture_timestamp,
        }


@dataclass
class ShapeDiff:
    """Difference between a shape before and after transformation."""
    shape_name: str
    field: str
    before: Any
    after: Any
    is_expected: bool = True
    note: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {
            "shape": self.shape_name,
            "field": self.field,
            "before": self.before,
            "after": self.after,
            "expected": self.is_expected,
            "note": self.note,
        }


@dataclass
class SlideDiff:
    """Structural diff for a single slide."""
    slide_num: int
    layout_type: str
    changes: List[ShapeDiff] = field(default_factory=list)
    anomalies: List[str] = field(default_factory=list)

    @property
    def change_count(self) -> int:
        return len(self.changes)

    @property
    def anomaly_count(self) -> int:
        return len(self.anomalies)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "slide_num": self.slide_num,
            "layout_type": self.layout_type,
            "change_count": self.change_count,
            "anomaly_count": self.anomaly_count,
            "changes": [c.to_dict() for c in self.changes],
            "anomalies": self.anomalies,
        }


@dataclass
class DeckDiff:
    """Complete structural diff for a presentation."""
    deck_name: str
    slide_diffs: List[SlideDiff] = field(default_factory=list)
    global_anomalies: List[str] = field(default_factory=list)

    @property
    def total_changes(self) -> int:
        return sum(s.change_count for s in self.slide_diffs)

    @property
    def total_anomalies(self) -> int:
        return sum(s.anomaly_count for s in self.slide_diffs) + len(self.global_anomalies)

    @property
    def slides_with_changes(self) -> int:
        return sum(1 for s in self.slide_diffs if s.change_count > 0)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "deck_name": self.deck_name,
            "total_changes": self.total_changes,
            "total_anomalies": self.total_anomalies,
            "slides_with_changes": self.slides_with_changes,
            "global_anomalies": self.global_anomalies,
            "slide_diffs": [s.to_dict() for s in self.slide_diffs],
        }


# ─────────────────────────────────────────────────────────────────────────────
# Snapshot Capture
# ─────────────────────────────────────────────────────────────────────────────

class StructureSnapshotter:
    """
    Captures structural snapshots of presentations.
    """

    # Shape type integer → human label mapping (python-pptx MSO_SHAPE_TYPE)
    SHAPE_TYPE_NAMES = {
        1: "AUTO_SHAPE",
        2: "CALLOUT",
        3: "CHART",
        4: "COMMENT",
        5: "FREEFORM",
        6: "GROUP",
        7: "EMBEDDED_OLE_OBJECT",
        8: "FORM_CONTROL",
        9: "LINE",
        10: "LINKED_OLE_OBJECT",
        11: "LINKED_PICTURE",
        12: "OLE_CONTROL_OBJECT",
        13: "PICTURE",
        14: "PLACEHOLDER",
        15: "SCRIPT_ANCHOR",
        16: "TABLE",
        17: "TEXT_BOX",
        18: "TEXT_EFFECT",
        19: "MEDIA",
        20: "WEB_VIDEO",
    }

    def capture(self, pptx_path: str) -> DeckSnapshot:
        """
        Capture a structural snapshot of a presentation.

        Args:
            pptx_path: Path to the .pptx file.

        Returns:
            DeckSnapshot with per-slide shape details.
        """
        prs = Presentation(pptx_path)
        snapshot = DeckSnapshot(
            file_path=pptx_path,
            slide_count=len(prs.slides),
            slide_width=int(prs.slide_width),
            slide_height=int(prs.slide_height),
        )

        for slide_idx, slide in enumerate(prs.slides):
            slide_snap = self._capture_slide(slide, slide_idx + 1)
            snapshot.slides.append(slide_snap)

        return snapshot

    def _capture_slide(self, slide, slide_num: int) -> SlideSnapshot:
        """Capture snapshot of a single slide."""
        layout = slide.slide_layout
        layout_name = getattr(layout, 'name', 'unknown')
        layout_type = self._get_layout_type(layout)

        slide_snap = SlideSnapshot(
            slide_num=slide_num,
            layout_name=layout_name,
            layout_type=layout_type,
            shape_count=len(slide.shapes),
        )

        for shape in slide.shapes:
            shape_snap = self._capture_shape(shape)
            slide_snap.shapes.append(shape_snap)

        return slide_snap

    def _capture_shape(self, shape) -> ShapeSnapshot:
        """Capture snapshot of a single shape."""
        shape_type_int = getattr(shape.shape_type, 'real', int(shape.shape_type))
        shape_type_name = self.SHAPE_TYPE_NAMES.get(shape_type_int, f"TYPE_{shape_type_int}")

        has_text = getattr(shape, 'has_text_frame', False)
        text_preview = ""
        rtl_flag = None
        alignment = None

        if has_text and shape.has_text_frame:
            try:
                full_text = shape.text_frame.text or ""
                text_preview = full_text[:60].replace('\n', ' ').strip()

                # Get RTL flag and alignment from first non-empty paragraph
                for para in shape.text_frame.paragraphs:
                    if not para.text.strip():
                        continue
                    try:
                        pPr = para._p.find(
                            '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr'
                        )
                        if pPr is not None:
                            rtl_val = pPr.get('rtl')
                            if rtl_val is not None:
                                rtl_flag = rtl_val in ('1', 'true')
                            algn_val = pPr.get('algn')
                            if algn_val is not None:
                                alignment = algn_val
                    except Exception:
                        pass
                    break
            except Exception:
                pass

        is_ph = getattr(shape, 'is_placeholder', False)
        ph_type = None
        ph_idx = None
        if is_ph:
            try:
                ph_fmt = shape.placeholder_format
                if ph_fmt:
                    ph_type = str(ph_fmt.type).split('.')[-1].lower()
                    ph_idx = ph_fmt.idx
            except Exception:
                pass

        return ShapeSnapshot(
            name=getattr(shape, 'name', ''),
            shape_type=shape_type_name,
            left=int(getattr(shape, 'left', 0) or 0),
            top=int(getattr(shape, 'top', 0) or 0),
            width=int(getattr(shape, 'width', 0) or 0),
            height=int(getattr(shape, 'height', 0) or 0),
            has_text=has_text,
            text_preview=text_preview,
            is_placeholder=is_ph,
            placeholder_type=ph_type,
            placeholder_idx=ph_idx,
            rtl_flag=rtl_flag,
            alignment=alignment,
        )

    def _get_layout_type(self, layout) -> str:
        """Extract the ST_SlideLayoutType string from a layout."""
        try:
            # Access the underlying XML
            cSld = layout._element
            # Look for the type attribute on the root or its parent
            # The layout type is usually on the <p:cSld> element's parent <p:sldLayout>
            sld_layout = cSld.getparent() if hasattr(cSld, 'getparent') else cSld
            type_attr = sld_layout.get('type', '')
            if type_attr:
                return type_attr
        except Exception:
            pass
        return getattr(layout, 'name', 'unknown')


# ─────────────────────────────────────────────────────────────────────────────
# Structural Diff Engine
# ─────────────────────────────────────────────────────────────────────────────

class StructuralDiffer:
    """
    Compares two DeckSnapshots (before and after transformation) and
    produces a structured diff.
    """

    # Tolerance for position/size comparison (in EMU)
    # ±200 000 EMU ≈ ±0.22" — allow small floating-point drift
    POSITION_TOLERANCE_EMU = 200_000

    def diff(self, before: DeckSnapshot, after: DeckSnapshot) -> DeckDiff:
        """
        Compare two snapshots and return the structural diff.

        Args:
            before: Pre-transformation snapshot.
            after:  Post-transformation snapshot.

        Returns:
            DeckDiff with per-slide change lists and anomalies.
        """
        deck_diff = DeckDiff(deck_name=Path(before.file_path).stem)

        # Global checks
        if before.slide_count != after.slide_count:
            deck_diff.global_anomalies.append(
                f"Slide count changed: {before.slide_count} → {after.slide_count}"
            )

        if before.slide_width != after.slide_width:
            deck_diff.global_anomalies.append(
                f"Slide width changed: {before.slide_width} → {after.slide_width}"
            )

        # Per-slide diff
        for i, (before_slide, after_slide) in enumerate(
            zip(before.slides, after.slides)
        ):
            slide_diff = self._diff_slide(before_slide, after_slide)
            deck_diff.slide_diffs.append(slide_diff)

        return deck_diff

    def _diff_slide(self, before: SlideSnapshot, after: SlideSnapshot) -> SlideDiff:
        """Diff a single slide."""
        slide_diff = SlideDiff(
            slide_num=before.slide_num,
            layout_type=before.layout_type,
        )

        if before.shape_count != after.shape_count:
            slide_diff.anomalies.append(
                f"Shape count changed: {before.shape_count} → {after.shape_count}"
            )

        # Match shapes by name
        before_shapes = {s.name: s for s in before.shapes}
        after_shapes = {s.name: s for s in after.shapes}

        for name, before_shape in before_shapes.items():
            if name not in after_shapes:
                slide_diff.anomalies.append(f"Shape disappeared: {name!r}")
                continue

            after_shape = after_shapes[name]
            changes = self._diff_shape(before_shape, after_shape)
            slide_diff.changes.extend(changes)

        for name in after_shapes:
            if name not in before_shapes:
                slide_diff.anomalies.append(f"New shape appeared: {name!r}")

        return slide_diff

    def _diff_shape(self, before: ShapeSnapshot,
                    after: ShapeSnapshot) -> List[ShapeDiff]:
        """Compare individual shape properties."""
        diffs = []

        # Position / size changes
        for field_name in ('left', 'top', 'width', 'height'):
            b_val = getattr(before, field_name)
            a_val = getattr(after, field_name)

            if abs(b_val - a_val) > self.POSITION_TOLERANCE_EMU:
                # Classify as expected or anomalous
                is_expected = self._is_expected_position_change(
                    field_name, before, after
                )
                diffs.append(ShapeDiff(
                    shape_name=before.name,
                    field=field_name,
                    before=b_val,
                    after=a_val,
                    is_expected=is_expected,
                    note=self._position_change_note(field_name, b_val, a_val),
                ))

        # RTL flag changes
        if before.rtl_flag != after.rtl_flag:
            is_expected = (before.rtl_flag is None or before.rtl_flag is False) and \
                          (after.rtl_flag is True)
            diffs.append(ShapeDiff(
                shape_name=before.name,
                field='rtl_flag',
                before=before.rtl_flag,
                after=after.rtl_flag,
                is_expected=is_expected,
                note='RTL direction set' if is_expected else 'Unexpected RTL change',
            ))

        # Alignment changes
        if before.alignment != after.alignment:
            is_expected = (after.alignment in ('r', 'ctr')) and before.has_text
            diffs.append(ShapeDiff(
                shape_name=before.name,
                field='alignment',
                before=before.alignment,
                after=after.alignment,
                is_expected=is_expected,
                note='Alignment set for RTL' if is_expected else 'Unexpected alignment change',
            ))

        return diffs

    def _is_expected_position_change(
        self, field: str, before: ShapeSnapshot, after: ShapeSnapshot
    ) -> bool:
        """
        Determine if a position change is expected for RTL transformation.

        Expected changes:
        - 'left': Shapes mirror X (expected for non-centered non-placeholder shapes)
        - 'top', 'width', 'height': Should NOT change in RTL transform
        """
        if field == 'left':
            # Left position changes are expected for mirrored shapes
            return True
        # Top/width/height changes are unexpected unless very small
        return False

    def _position_change_note(self, field: str, before: int, after: int) -> str:
        """Generate a human-readable note for a position change."""
        delta = after - before
        delta_in = delta / 914400  # EMU to inches
        return f"Δ{field}={delta:+d} EMU ({delta_in:+.2f}\")"


# ─────────────────────────────────────────────────────────────────────────────
# HTML Report Generator
# ─────────────────────────────────────────────────────────────────────────────

class HTMLReportGenerator:
    """
    Generates an HTML report from a DeckDiff.
    """

    CSS = """
<style>
body { font-family: 'Segoe UI', Arial, sans-serif; margin: 20px; color: #222; }
h1 { color: #1a1a2e; }
h2 { color: #16213e; border-bottom: 2px solid #eee; padding-bottom: 5px; }
.summary-box { background: #f8f9fa; border: 1px solid #dee2e6; border-radius: 8px;
               padding: 15px; margin: 15px 0; display: inline-block; min-width: 200px; }
.stat-number { font-size: 2em; font-weight: bold; color: #0d6efd; }
.stat-label { color: #666; font-size: 0.85em; }
.slide-section { margin: 20px 0; }
.slide-header { background: #e9ecef; padding: 8px 15px; border-radius: 5px;
                cursor: pointer; user-select: none; }
.slide-header:hover { background: #dee2e6; }
.slide-content { padding: 10px 15px; display: none; border: 1px solid #dee2e6;
                 border-top: none; border-radius: 0 0 5px 5px; }
.change-table { width: 100%; border-collapse: collapse; margin: 10px 0; }
.change-table th { background: #343a40; color: white; padding: 6px 10px; text-align: left; }
.change-table td { padding: 5px 10px; border-bottom: 1px solid #dee2e6; font-size: 0.85em; }
.change-table tr:hover td { background: #f8f9fa; }
.expected { color: #198754; }
.unexpected { color: #dc3545; font-weight: bold; }
.anomaly { background: #fff3cd; padding: 5px 10px; border-left: 3px solid #ffc107;
           margin: 5px 0; font-size: 0.85em; }
.global-anomaly { background: #f8d7da; padding: 8px 12px; border-left: 4px solid #dc3545;
                  margin: 8px 0; }
.badge { display: inline-block; padding: 2px 8px; border-radius: 10px;
         font-size: 0.75em; font-weight: bold; }
.badge-changes { background: #cfe2ff; color: #084298; }
.badge-anomalies { background: #f8d7da; color: #842029; }
.badge-clean { background: #d1e7dd; color: #0a3622; }
toggle-btn { background: none; border: none; font-size: 1em; cursor: pointer; }
</style>
<script>
function toggleSlide(id) {
    var el = document.getElementById(id);
    if (el) el.style.display = el.style.display === 'block' ? 'none' : 'block';
}
function expandAll() {
    document.querySelectorAll('.slide-content').forEach(e => e.style.display = 'block');
}
function collapseAll() {
    document.querySelectorAll('.slide-content').forEach(e => e.style.display = 'none');
}
</script>
"""

    def generate(self, diff: DeckDiff, before: DeckSnapshot,
                 after: DeckSnapshot) -> str:
        """
        Generate an HTML report comparing before/after structural states.

        Args:
            diff: DeckDiff produced by StructuralDiffer.
            before: Pre-transformation snapshot.
            after: Post-transformation snapshot.

        Returns:
            HTML string.
        """
        html = []
        html.append("<!DOCTYPE html><html lang='en'><head>")
        html.append("<meta charset='UTF-8'>")
        html.append(f"<title>SlideShift v2 — {diff.deck_name}</title>")
        html.append(self.CSS)
        html.append("</head><body>")

        # Header
        html.append(f"<h1>SlideShift v2 Structural Report: {diff.deck_name}</h1>")
        html.append(f"<p><strong>Input:</strong> {before.file_path}</p>")

        # Summary boxes
        html.append("<div style='display:flex; gap:20px; flex-wrap:wrap;'>")
        html.append(self._summary_box("Slides", before.slide_count, ""))
        html.append(self._summary_box("Total Changes", diff.total_changes, ""))
        html.append(self._summary_box("Anomalies", diff.total_anomalies, ""))
        html.append(
            self._summary_box("Slides Changed",
                              f"{diff.slides_with_changes}/{before.slide_count}", "")
        )
        html.append("</div>")

        # Global anomalies
        if diff.global_anomalies:
            html.append("<h2>Global Anomalies</h2>")
            for anomaly in diff.global_anomalies:
                html.append(f"<div class='global-anomaly'>{anomaly}</div>")

        # Per-slide sections
        html.append("<h2>Per-Slide Analysis</h2>")
        html.append(
            "<button onclick='expandAll()'>Expand All</button> "
            "<button onclick='collapseAll()'>Collapse All</button>"
        )

        for slide_diff in diff.slide_diffs:
            s_before = before.slides[slide_diff.slide_num - 1]
            s_after = after.slides[slide_diff.slide_num - 1]
            html.append(self._slide_section(slide_diff, s_before, s_after))

        html.append("</body></html>")
        return "\n".join(html)

    def _summary_box(self, label: str, value: Any, unit: str) -> str:
        return (
            f"<div class='summary-box'>"
            f"<div class='stat-number'>{value}{unit}</div>"
            f"<div class='stat-label'>{label}</div>"
            f"</div>"
        )

    def _slide_section(self, slide_diff: SlideDiff,
                       s_before: SlideSnapshot, s_after: SlideSnapshot) -> str:
        html = []
        slide_id = f"slide_{slide_diff.slide_num}"

        # Badge
        if slide_diff.anomaly_count > 0:
            badge = f"<span class='badge badge-anomalies'>{slide_diff.anomaly_count} anomalies</span>"
        elif slide_diff.change_count > 0:
            badge = f"<span class='badge badge-changes'>{slide_diff.change_count} changes</span>"
        else:
            badge = "<span class='badge badge-clean'>clean</span>"

        html.append("<div class='slide-section'>")
        html.append(
            f"<div class='slide-header' onclick=\"toggleSlide('{slide_id}')\">"
            f"Slide {slide_diff.slide_num} — {s_before.layout_name} "
            f"({slide_diff.layout_type}) {badge}"
            f"</div>"
        )
        html.append(f"<div class='slide-content' id='{slide_id}'>")

        # Anomalies
        for anomaly in slide_diff.anomalies:
            html.append(f"<div class='anomaly'>⚠️ {anomaly}</div>")

        # Changes table
        if slide_diff.changes:
            html.append("<table class='change-table'>")
            html.append(
                "<tr><th>Shape</th><th>Field</th><th>Before</th>"
                "<th>After</th><th>Status</th><th>Note</th></tr>"
            )
            for change in slide_diff.changes:
                css_class = 'expected' if change.is_expected else 'unexpected'
                status_text = '✓ Expected' if change.is_expected else '✗ Unexpected'
                html.append(
                    f"<tr>"
                    f"<td>{change.shape_name}</td>"
                    f"<td>{change.field}</td>"
                    f"<td>{change.before}</td>"
                    f"<td>{change.after}</td>"
                    f"<td class='{css_class}'>{status_text}</td>"
                    f"<td>{change.note}</td>"
                    f"</tr>"
                )
            html.append("</table>")
        else:
            html.append("<p><em>No structural changes on this slide.</em></p>")

        html.append("</div></div>")
        return "\n".join(html)


# ─────────────────────────────────────────────────────────────────────────────
# Test Runner
# ─────────────────────────────────────────────────────────────────────────────

class TestHarness:
    """
    Orchestrates the full test run for one or more PPTX files.
    """

    def __init__(self, output_dir: str = ".", verbose: bool = False):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.verbose = verbose
        self.snapshotter = StructureSnapshotter()
        self.differ = StructuralDiffer()
        self.reporter = HTMLReportGenerator()
        self.run_results: List[Dict[str, Any]] = []

    def run_single(self, pptx_path: str) -> Dict[str, Any]:
        """
        Run the full test pipeline on a single PPTX file.

        Returns a result dict with stats and output paths.
        """
        deck_name = Path(pptx_path).stem
        start_time = time.monotonic()
        result = {
            "deck": deck_name,
            "input": pptx_path,
            "success": False,
            "error": None,
            "changes": 0,
            "anomalies": 0,
            "output_pptx": None,
            "report_html": None,
            "elapsed_ms": 0,
        }

        try:
            logger.info("[%s] Capturing pre-transform snapshot...", deck_name)
            before_snap = self.snapshotter.capture(pptx_path)

            # Copy the file to a temp location for transformation
            import shutil
            import tempfile
            with tempfile.NamedTemporaryFile(
                suffix='.pptx', delete=False
            ) as tmp:
                tmp_path = tmp.name

            shutil.copy2(pptx_path, tmp_path)

            # Load the copy and transform it
            logger.info("[%s] Applying RTL transformation...", deck_name)
            prs = Presentation(tmp_path)

            try:
                registry = TemplateRegistry(
                    int(prs.slide_width), int(prs.slide_height)
                )
                master_transformer = MasterLayoutTransformer(prs, registry)
                master_transformer.transform_all_masters()
                master_transformer.transform_all_layouts()

                # Classify slides then transform content
                analyzer = LayoutAnalyzer()
                layout_classifications = analyzer.classify_slides(prs)
                slide_transformer = SlideContentTransformer(
                    prs, layout_classifications, {}
                )
                slide_transformer.transform_all_slides()
            except NameError as e:
                logger.warning("[%s] Transform modules not available: %s", deck_name, e)
                logger.warning("[%s] Using no-op transform for structural test", deck_name)

            # Save transformed output
            output_pptx = str(self.output_dir / f"{deck_name}_rtl.pptx")
            prs.save(output_pptx)
            result["output_pptx"] = output_pptx
            logger.info("[%s] Saved RTL output to %s", deck_name, output_pptx)

            # Capture post-transform snapshot
            logger.info("[%s] Capturing post-transform snapshot...", deck_name)
            after_snap = self.snapshotter.capture(output_pptx)

            # Diff
            diff = self.differ.diff(before_snap, after_snap)
            result["changes"] = diff.total_changes
            result["anomalies"] = diff.total_anomalies

            # Generate HTML report
            html_content = self.reporter.generate(diff, before_snap, after_snap)
            report_path = str(self.output_dir / f"{deck_name}_report.html")
            with open(report_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            result["report_html"] = report_path
            logger.info("[%s] Report written to %s", deck_name, report_path)

            # Save JSON diff
            json_path = str(self.output_dir / f"{deck_name}_diff.json")
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(diff.to_dict(), f, indent=2)

            result["success"] = True

            # Cleanup temp file
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

        except Exception as exc:
            logger.error("[%s] Error: %s", deck_name, exc, exc_info=True)
            result["error"] = str(exc)

        result["elapsed_ms"] = (time.monotonic() - start_time) * 1000
        self.run_results.append(result)
        return result

    def run_all(self, pptx_paths: List[str]) -> Dict[str, Any]:
        """
        Run the test pipeline on multiple PPTX files.

        Returns an aggregate summary.
        """
        for path in pptx_paths:
            logger.info("Processing: %s", path)
            result = self.run_single(path)
            status = "✓" if result["success"] else "✗"
            print(
                f"  {status} {Path(path).name}: "
                f"{result['changes']} changes, {result['anomalies']} anomalies "
                f"({result['elapsed_ms']:.0f}ms)"
            )

        summary = self._build_summary()
        return summary

    def _build_summary(self) -> Dict[str, Any]:
        """Build an aggregate summary of all run results."""
        total = len(self.run_results)
        successful = sum(1 for r in self.run_results if r["success"])
        total_changes = sum(r["changes"] for r in self.run_results)
        total_anomalies = sum(r["anomalies"] for r in self.run_results)
        total_elapsed = sum(r["elapsed_ms"] for r in self.run_results)

        return {
            "total_decks": total,
            "successful": successful,
            "failed": total - successful,
            "total_changes": total_changes,
            "total_anomalies": total_anomalies,
            "total_elapsed_ms": total_elapsed,
            "per_deck": self.run_results,
        }

    def save_run_summary(self, output_path: str) -> None:
        """Save the run summary to a JSON file."""
        summary = self._build_summary()
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(summary, f, indent=2)
        logger.info("Run summary saved to %s", output_path)


# ─────────────────────────────────────────────────────────────────────────────
# CLI Entry Point
# ─────────────────────────────────────────────────────────────────────────────

def find_pptx_files(directory: str) -> List[str]:
    """Find all .pptx files in a directory (non-recursive)."""
    dir_path = Path(directory)
    return sorted(
        str(p) for p in dir_path.iterdir()
        if p.suffix.lower() == '.pptx' and not p.name.startswith('~')
    )


def main():
    parser = argparse.ArgumentParser(
        description='SlideShift v2 Structural Test Harness',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python test_harness.py deck.pptx
  python test_harness.py deck1.pptx deck2.pptx --output-dir reports/
  python test_harness.py --dir /path/to/decks/ --output-dir reports/
""",
    )
    parser.add_argument(
        'files', nargs='*', metavar='FILE',
        help='PPTX file(s) to process',
    )
    parser.add_argument(
        '--dir', metavar='DIRECTORY',
        help='Directory of PPTX files to process',
    )
    parser.add_argument(
        '--output-dir', default='.', metavar='DIR',
        help='Output directory for reports (default: current dir)',
    )
    parser.add_argument(
        '--verbose', '-v', action='store_true',
        help='Enable verbose logging',
    )
    parser.add_argument(
        '--log-file', metavar='FILE',
        help='Write logs to file',
    )

    args = parser.parse_args()

    # Configure logging
    log_level = logging.DEBUG if args.verbose else logging.INFO
    handlers = [logging.StreamHandler()]
    if args.log_file:
        handlers.append(logging.FileHandler(args.log_file))
    logging.basicConfig(
        level=log_level,
        format='%(asctime)s %(levelname)s %(name)s: %(message)s',
        handlers=handlers,
    )

    # Collect input files
    pptx_files = list(args.files)
    if args.dir:
        pptx_files.extend(find_pptx_files(args.dir))

    if not pptx_files:
        print("ERROR: No PPTX files specified. Use --help for usage.")
        sys.exit(1)

    # Validate inputs
    for f in pptx_files:
        if not Path(f).exists():
            print(f"ERROR: File not found: {f}")
            sys.exit(1)

    print(f"SlideShift v2 Test Harness")
    print(f"Processing {len(pptx_files)} deck(s)...")
    print()

    # Run
    harness = TestHarness(
        output_dir=args.output_dir,
        verbose=args.verbose,
    )
    summary = harness.run_all(pptx_files)

    # Save run summary
    run_summary_path = str(Path(args.output_dir) / 'run_summary.json')
    harness.save_run_summary(run_summary_path)

    # Print summary
    print()
    print("═" * 50)
    print("Run Summary")
    print("═" * 50)
    print(f"  Decks processed : {summary['total_decks']}")
    print(f"  Successful       : {summary['successful']}")
    print(f"  Failed           : {summary['failed']}")
    print(f"  Total changes    : {summary['total_changes']}")
    print(f"  Total anomalies  : {summary['total_anomalies']}")
    print(f"  Total time       : {summary['total_elapsed_ms']:.0f}ms")
    print()
    print(f"  Run summary → {run_summary_path}")


if __name__ == "__main__":
    main()
