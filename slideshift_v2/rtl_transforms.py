"""
rtl_transforms.py — Deterministic RTL transformation functions.

SlideShift v2: Template-First Deterministic RTL Transformation Engine.

Phase 2: MasterLayoutTransformer  — transforms slide masters and layouts FIRST
Phase 3: SlideContentTransformer  — transforms individual content slides

Design principles (from 4-model architectural consensus):
1. Masters are style sheets — set direction/language defaults, mirror logos only.
   NEVER set algn at master/layout level (context-sensitive, slide-level concern).
   NEVER apply flipH (corrupts logos and gradients).
2. Layouts own placeholder geometry — mirror and swap placeholder positions.
3. Slides own content — placeholders inherit from transformed layout by deleting
   local position overrides; freeform shapes are mirrored explicitly.
4. Alignment (algn) is ALWAYS written explicitly at the paragraph level on slides.
"""

from __future__ import annotations

import logging
from copy import deepcopy
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple

from lxml import etree

from .utils import (
    A_NS, P_NS, R_NS,
    mirror_x,
    swap_positions,
    has_arabic,
    compute_script_ratio,
    qn,
    ensure_pPr,
    set_rtl_on_paragraph,
    set_alignment_on_paragraph,
    get_placeholder_info,
    get_placeholder_info_from_xml,
    set_body_pr_rtl_col,
    set_defRPr_lang,
    iter_paragraphs,
    iter_runs,
    get_run_text,
    bounds_check_emu,
    clamp_emu,
)

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Data classes
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class TransformReport:
    """Summary report produced by each transformation phase."""
    phase: str  # 'master', 'layout', 'slide', 'typography'
    total_changes: int = 0
    changes_by_type: Dict[str, int] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)

    def add(self, change_type: str, count: int = 1) -> None:
        """Increment a change-type counter and the total."""
        self.changes_by_type[change_type] = (
            self.changes_by_type.get(change_type, 0) + count
        )
        self.total_changes += count

    def warn(self, msg: str) -> None:
        self.warnings.append(msg)
        logger.warning('[%s] %s', self.phase, msg)

    def error(self, msg: str) -> None:
        self.errors.append(msg)
        logger.error('[%s] %s', self.phase, msg)

    def merge(self, other: 'TransformReport') -> None:
        """Merge another report's counters into this one."""
        self.total_changes += other.total_changes
        for k, v in other.changes_by_type.items():
            self.changes_by_type[k] = self.changes_by_type.get(k, 0) + v
        self.warnings.extend(other.warnings)
        self.errors.extend(other.errors)


# ─────────────────────────────────────────────────────────────────────────────
# Internal constants
# ─────────────────────────────────────────────────────────────────────────────

# Logo detection: image width must be less than this fraction of slide width
_LOGO_MAX_WIDTH_FRACTION = 0.20

# Position tolerance: skip mirror if the change would be smaller than this (EMU)
_POSITION_TOLERANCE_EMU = 50_000  # ≈ 0.055 inches

# Placeholder type strings that are "title-like" (keep centered or right-align)
_TITLE_PH_TYPES = frozenset({'title', 'ctrTitle', 'center_title'})

# Placeholder type strings that should always be left-aligned (footers, dates)
_FOOTER_PH_TYPES = frozenset({'ftr', 'sldNum', 'dt', 'footer', 'slideNumber', 'date_time'})


# ─────────────────────────────────────────────────────────────────────────────
# MasterLayoutTransformer — Phase 2
# ─────────────────────────────────────────────────────────────────────────────

class MasterLayoutTransformer:
    """
    Transforms slide masters and layouts for RTL (Phase 2).

    Operates BEFORE any content slides are touched so that placeholders
    on slides can inherit the correct RTL positions from their layouts.

    What this does:
    - Masters: set RTL text-direction defaults (rtl, rtlCol, lang); mirror logos.
    - Layouts: mirror/swap placeholder X positions; set RTL defaults.

    What this does NOT do (by design):
    - Set algn on masters/layouts (context-sensitive; set at slide paragraph level).
    - Apply flipH to any shape (corrupts logos, inverts brand gradients).
    - Modify shape heights or fonts (typography phase).
    """

    def __init__(self, presentation, template_registry=None):
        """
        Args:
            presentation: python-pptx Presentation object.
            template_registry: Optional TemplateRegistry instance with layout rules.
                When None, a built-in default ruleset is used.
        """
        self.prs = presentation
        self.template_registry = template_registry
        self._slide_width = int(presentation.slide_width)
        self._slide_height = int(presentation.slide_height)

    # ───────────────────────────────────────────────────────────────────
    # Public entry points
    # ───────────────────────────────────────────────────────────────────

    def transform_all_masters(self) -> TransformReport:
        """
        Transform all slide masters for RTL.

        Returns:
            TransformReport summarising all changes made.
        """
        report = TransformReport(phase='master')
        for idx, master in enumerate(self.prs.slide_masters):
            try:
                count = self._transform_master(master)
                report.add('master_transformed', count)
            except Exception as exc:
                report.error(f'master[{idx}]: {exc}')
        return report

    def transform_all_layouts(self) -> TransformReport:
        """
        Transform all slide layouts for RTL.

        Returns:
            TransformReport summarising all changes made.
        """
        report = TransformReport(phase='layout')
        for m_idx, master in enumerate(self.prs.slide_masters):
            for l_idx, layout in enumerate(master.slide_layouts):
                try:
                    count = self._transform_layout(layout)
                    report.add('layout_transformed', count)
                except Exception as exc:
                    report.error(f'master[{m_idx}].layout[{l_idx}]: {exc}')
        return report

    # ───────────────────────────────────────────────────────────────────
    # Master transformation
    # ───────────────────────────────────────────────────────────────────

    def _transform_master(self, master) -> int:
        """
        Transform a single slide master.

        Returns:
            Count of XML attribute writes performed.
        """
        changes = 0
        xml_el = master._element

        # 1. RTL text-direction defaults on all text body elements
        changes += self._apply_rtl_direction_defaults(xml_el)

        # 2. Arabic language defaults on all defRPr elements
        changes += self._apply_arabic_lang_defaults(xml_el)

        # 3. Set RTL in master txStyles (paragraph level defaults per style)
        changes += self._set_master_text_styles_rtl(master)

        # 4. Mirror logo images (position only — NO flipH)
        changes += self._mirror_logo_images(master)

        # 5. Mirror small brand text elements (company names, lettermarks)
        changes += self._mirror_brand_elements(master)

        return changes

    def _apply_rtl_direction_defaults(self, xml_element) -> int:
        """
        Set direction defaults on master/layout elements.
        Only sets rtlCol='1' on bodyPr (safe). Does NOT set rtl='1' on lstStyle
        or defPPr (that would corrupt English text).
        """
        changes = 0
        for body_pr in xml_element.iter(f'{{{A_NS}}}bodyPr'):
            body_pr.set('rtlCol', '1')
            changes += 1
        return changes

    def _apply_arabic_lang_defaults(self, xml_element) -> int:
        """Set lang='ar-SA' on all <a:defRPr> elements."""
        changes = 0
        for def_rPr in xml_element.iter(f'{{{A_NS}}}defRPr'):
            def_rPr.set('lang', 'ar-SA')
            changes += 1
        return changes

    def _set_master_text_styles_rtl(self, master) -> int:
        """
        Set RTL-related defaults in the master's txStyles element.
        Only sets language defaults on defRPr. RTL direction handled per-paragraph at slide level.
        """
        changes = 0
        try:
            xml_el = master._element
            tx_styles = xml_el.find(f'{{{P_NS}}}txStyles')
            if tx_styles is None:
                return 0

            for style_name in ('titleStyle', 'bodyStyle', 'otherStyle'):
                style_elem = tx_styles.find(f'{{{A_NS}}}{style_name}')
                if style_elem is None:
                    continue
                for level in range(1, 10):
                    for lvl_pPr in style_elem.findall(f'{{{A_NS}}}lvl{level}pPr'):
                        defRPr = lvl_pPr.find(f'{{{A_NS}}}defRPr')
                        if defRPr is not None:
                            defRPr.set('lang', 'ar-SA')
                            changes += 1
        except Exception as exc:
            logger.warning('_set_master_text_styles_rtl: %s', exc)
        return changes

    def _mirror_brand_elements(self, element) -> int:
        """
        Mirror small non-placeholder text shapes on masters/layouts.
        Handles rotated shapes and negative positions.
        """
        mirrored = 0
        try:
            for shape in element.shapes:
                sp_el = shape._element
                ph = sp_el.find(f'.//{{{P_NS}}}ph')
                if ph is not None:
                    continue

                if not (getattr(shape, 'has_text_frame', False) and shape.has_text_frame):
                    continue
                text = shape.text_frame.text or ''
                if not text.strip():
                    continue

                left = shape.left
                width = shape.width
                height = shape.height
                if left is None or width is None:
                    continue

                rotation_deg = 0
                for xfrm in sp_el.iter(f'{{{A_NS}}}xfrm'):
                    rot_val = xfrm.get('rot', '0')
                    try:
                        rotation_deg = int(rot_val) / 60000
                    except (ValueError, TypeError):
                        pass
                    break

                is_rotated_90 = abs(rotation_deg) in (90, 270)
                if is_rotated_90:
                    visual_width = height if height else width
                else:
                    visual_width = width

                if visual_width > self._slide_width * 0.30:
                    continue

                new_left = mirror_x(left, width, self._slide_width)

                if left < 0:
                    bleed = abs(left)
                    visible_right_edge = left + width
                    if visible_right_edge <= 0:
                        continue
                    new_left = self._slide_width - width + bleed

                if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                    continue

                shape.left = new_left
                mirrored += 1

        except Exception as exc:
            logger.warning('_mirror_brand_elements: %s', exc)

        return mirrored

    def _mirror_logo_images(self, element) -> int:
        """
        Mirror the horizontal position of small logo images on a master or layout.
        Strict whitelist: small picture element, no text frame, width <20%, not placeholder, has blipFill.
        Applies position translation ONLY — no flipH.
        """
        mirrored = 0
        try:
            for shape in element.shapes:
                if not self._is_logo_shape(shape, self._slide_width):
                    continue

                original_left = shape.left
                original_width = shape.width
                if original_left is None or original_width is None:
                    continue

                new_left = mirror_x(original_left, original_width, self._slide_width)

                if not bounds_check_emu(new_left, self._slide_width):
                    continue

                if abs(new_left - original_left) < _POSITION_TOLERANCE_EMU:
                    continue

                shape.left = new_left
                mirrored += 1

        except Exception as exc:
            logger.warning('_mirror_logo_images: %s', exc)

        return mirrored

    def _is_logo_shape(self, shape, slide_width: int) -> bool:
        """
        Detect if a shape is likely a logo (small image in a corner of the slide).
        Returns True only if ALL conditions are met:
        1. Is a picture element (<p:pic>)
        2. Has no text frame (pure image)
        3. Width < 20% of slide width
        4. Is not a placeholder
        5. Has an actual image relationship (blipFill with rId embed)
        """
        try:
            sp_el = shape._element
            tag = sp_el.tag

            if not (tag.endswith('}pic') or tag == 'pic'):
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if getattr(shape, 'shape_type', None) != MSO_SHAPE_TYPE.PICTURE:
                        return False
                except ImportError:
                    return False

            if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
                return False

            shape_width = getattr(shape, 'width', None)
            if shape_width is None or shape_width >= slide_width * _LOGO_MAX_WIDTH_FRACTION:
                return False

            nv_pic_pr = sp_el.find(f'{{{P_NS}}}nvPicPr')
            if nv_pic_pr is not None:
                ph = nv_pic_pr.find(f'.//{{{P_NS}}}ph')
                if ph is not None:
                    return False

            blip_fill = sp_el.find(f'{{{P_NS}}}blipFill')
            if blip_fill is None:
                blip_fill = sp_el.find(f'{{{A_NS}}}blipFill')
            if blip_fill is not None:
                blip = blip_fill.find(f'{{{A_NS}}}blip')
                if blip is not None and blip.get(f'{{{R_NS}}}embed'):
                    return True

            return False
        except Exception:
            return False

    # ───────────────────────────────────────────────────────────────────
    # Layout transformation
    # ───────────────────────────────────────────────────────────────────

    def _transform_layout(self, layout) -> int:
        """
        Transform a single slide layout for RTL.
        Returns count of XML attribute writes performed.
        """
        changes = 0
        xml_el = layout._element
        layout_type = xml_el.get('type', 'cust')

        changes += self._apply_rtl_direction_defaults(xml_el)
        changes += self._apply_arabic_lang_defaults(xml_el)
        changes += self._mirror_layout_placeholders(layout, layout_type)

        return changes

    def _mirror_layout_placeholders(self, layout, layout_type: str) -> int:
        """
        Mirror placeholder positions in a layout according to RTL rules.
        For two-column layouts, the two content placeholders are *swapped*.
        For all other placeholders, positions are *mirrored*.
        """
        changes = 0
        slide_width = self._slide_width

        two_column_types = frozenset({
            'twoColTx', 'twoObj', 'twoTxTwoObj',
            'txAndChart', 'chartAndTx', 'picTx',
        })

        if layout_type in two_column_types:
            changes += self._swap_two_column_placeholders(layout, slide_width)
        else:
            for shape in layout.placeholders:
                try:
                    left = shape.left
                    width = shape.width
                    if left is None or width is None:
                        continue
                    new_left = mirror_x(left, width, slide_width)
                    if not bounds_check_emu(new_left, slide_width):
                        continue
                    if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                        continue
                    shape.left = new_left
                    changes += 1
                except Exception as exc:
                    logger.debug('_mirror_layout_placeholders: %s', exc)

        return changes

    def _swap_two_column_placeholders(self, layout, slide_width: int) -> int:
        """
        Swap the horizontal positions of the two content-area placeholders
        in a two-column layout.
        """
        changes = 0
        content_placeholders = []
        title_placeholders = []

        for shape in layout.placeholders:
            ph_info = get_placeholder_info(shape)
            if ph_info is None:
                continue
            ph_type, ph_idx = ph_info
            if ph_type in _TITLE_PH_TYPES or ph_idx == 0:
                title_placeholders.append(shape)
            else:
                content_placeholders.append(shape)

        for shape in title_placeholders:
            try:
                left, width = shape.left, shape.width
                if left is None or width is None:
                    continue
                new_left = mirror_x(left, width, slide_width)
                if bounds_check_emu(new_left, slide_width) and abs(new_left - left) >= _POSITION_TOLERANCE_EMU:
                    shape.left = new_left
                    changes += 1
            except Exception as exc:
                logger.debug('title mirror: %s', exc)

        content_placeholders.sort(key=lambda s: getattr(s, 'left', 0) or 0)

        if len(content_placeholders) >= 2:
            left_ph = content_placeholders[0]
            right_ph = content_placeholders[-1]
            try:
                new_x_left, new_x_right = swap_positions(
                    left_ph.left, left_ph.width,
                    right_ph.left, right_ph.width,
                    slide_width,
                )
                left_ph.left = clamp_emu(new_x_left, slide_width)
                right_ph.left = clamp_emu(new_x_right, slide_width)
                changes += 2
            except Exception as exc:
                logger.warning('_swap_two_column_placeholders: %s', exc)

            for shape in content_placeholders[1:-1]:
                try:
                    left, width = shape.left, shape.width
                    if left is None or width is None:
                        continue
                    new_left = mirror_x(left, width, slide_width)
                    if bounds_check_emu(new_left, slide_width):
                        shape.left = new_left
                        changes += 1
                except Exception as exc:
                    logger.debug('extra content mirror: %s', exc)

        elif len(content_placeholders) == 1:
            shape = content_placeholders[0]
            try:
                left, width = shape.left, shape.width
                if left is not None and width is not None:
                    new_left = mirror_x(left, width, slide_width)
                    if bounds_check_emu(new_left, slide_width):
                        shape.left = new_left
                        changes += 1
            except Exception as exc:
                logger.debug('single content mirror: %s', exc)

        return changes

    def _mirror_shape_position(self, shape, slide_width_emu: int) -> bool:
        """Mirror a shape's X position using new_x = slide_width - (left + width)."""
        try:
            left = shape.left
            width = shape.width
            if left is None or width is None:
                return False
            new_left = mirror_x(left, width, slide_width_emu)
            if not bounds_check_emu(new_left, slide_width_emu):
                return False
            if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                return False
            shape.left = new_left
            return True
        except Exception as exc:
            logger.debug('_mirror_shape_position: %s', exc)
            return False

    def _swap_column_placeholders(
        self, layout, left_idx: int, right_idx: int, slide_width_emu: int
    ) -> bool:
        """Swap two placeholders' X positions given their placeholder indices."""
        left_ph = None
        right_ph = None
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == left_idx:
                left_ph = ph
            elif ph.placeholder_format.idx == right_idx:
                right_ph = ph

        if left_ph is None or right_ph is None:
            return False

        try:
            new_x_left, new_x_right = swap_positions(
                left_ph.left, left_ph.width,
                right_ph.left, right_ph.width,
                slide_width_emu,
            )
            left_ph.left = clamp_emu(new_x_left, slide_width_emu)
            right_ph.left = clamp_emu(new_x_right, slide_width_emu)
            return True
        except Exception as exc:
            logger.warning('_swap_column_placeholders(%d, %d): %s', left_idx, right_idx, exc)
            return False


# ─────────────────────────────────────────────────────────────────────────────
# SlideContentTransformer — Phase 3
# ─────────────────────────────────────────────────────────────────────────────

# NOTE: SlideContentTransformer (Phase 3) is defined in the full version of this file.
# This commit contains TransformReport and MasterLayoutTransformer only.
# The complete SlideContentTransformer with all 19 RTL fixes will follow in the next commit.
# See the original source at /home/user/workspace/slideshift_v2/rtl_transforms.py
