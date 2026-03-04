import logging
import time
from pathlib import Path
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Callable, Any
from pptx import Presentation

logger = logging.getLogger(__name__)

try:
    from slideshift_v2.property_resolver import PropertyResolver
    from slideshift_v2.layout_analyzer import LayoutAnalyzer
    from slideshift_v2.template_registry import TemplateRegistry
    from slideshift_v2.rtl_transforms import MasterLayoutTransformer, SlideContentTransformer
    from slideshift_v2.typography import TypographyNormalizer
    from slideshift_v2.structural_validator import StructuralValidator
    from slideshift_v2.models import (
        ResolvedPresentation, TransformReport, ValidationReport,
        PipelineConfig, PipelineResult
    )
except ImportError as e:
    logger.warning(f"Some v2 modules not yet available: {e}")
    
    @dataclass
    class PipelineConfig:
        input_path: str
        output_path: str
        translate_fn: Optional[Callable[[List[str]], Dict[str, str]]] = None
        skip_translation: bool = False
        max_font_reduction_pct: float = 20.0
        log_level: str = 'INFO'
        enable_telemetry: bool = False
        
    @dataclass
    class PipelineResult:
        success: bool
        output_path: Optional[str]
        phase_reports: Dict[str, Any]
        validation_report: Optional[Any]
        total_duration_ms: float
        error: Optional[str] = None


class SlideShiftV2Pipeline:
    """
    Template-first, deterministic RTL transformation pipeline.
    Phase 0: Parse & Resolve
    Phase 1: Translate
    Phase 2: Master & Layout Transform
    Phase 3: Slide Content Transform
    Phase 4: Typography Normalization
    Phase 5: Structural Validation
    """
    
    def __init__(self, config: 'PipelineConfig'):
        self.config = config
        self._phase_reports = {}
        numeric_level = getattr(logging, config.log_level.upper(), logging.INFO)
        logger.setLevel(numeric_level)
        
    def run(self) -> 'PipelineResult':
        pipeline_start = time.monotonic()
        logger.info(f"Starting SlideShift v2 pipeline for {self.config.input_path}")
        
        try:
            try:
                prs = Presentation(self.config.input_path)
            except Exception as e:
                raise ValueError(f"Failed to load presentation: {e}")
            
            resolved_prs = self._phase_0_resolve(prs)
            translation_map = self._phase_1_translate(resolved_prs)
            p2_report = self._phase_2_transform_masters_layouts(prs, resolved_prs)
            p3_report = self._phase_3_transform_slides(prs, resolved_prs, translation_map)
            p4_report = self._phase_4_typography(prs)
            val_report = self._phase_5_validate(prs, resolved_prs)
            
            try:
                Path(self.config.output_path).parent.mkdir(parents=True, exist_ok=True)
                prs.save(self.config.output_path)
            except Exception as e:
                raise IOError(f"Failed to save presentation: {e}")
            
            total_duration = (time.monotonic() - pipeline_start) * 1000
            logger.info(f"Pipeline completed successfully in {total_duration:.0f}ms")
            
            return PipelineResult(
                success=True,
                output_path=self.config.output_path,
                phase_reports=self._phase_reports,
                validation_report=val_report,
                total_duration_ms=total_duration
            )
            
        except Exception as e:
            logger.error(f"Pipeline failed: {e}", exc_info=True)
            total_duration = (time.monotonic() - pipeline_start) * 1000
            return PipelineResult(
                success=False,
                output_path=None,
                phase_reports=self._phase_reports,
                validation_report=None,
                total_duration_ms=total_duration,
                error=str(e)
            )
        
    def _phase_0_resolve(self, prs):
        start_time = time.monotonic()
        try:
            resolver = PropertyResolver(prs)
            resolved_prs = resolver.resolve_presentation()
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_0_resolve', duration, {"status": "success", "slides_resolved": len(prs.slides)})
            return resolved_prs
        except NameError as e:
            logger.warning(f"Phase 0 stubbed: {e}")
            self._log_phase('phase_0_resolve', 0, {"status": "stubbed"})
            return None
            
    def _phase_1_translate(self, resolved):
        start_time = time.monotonic()
        if self.config.skip_translation:
            self._log_phase('phase_1_translate', 0, {"status": "skipped"})
            return {}
        if not self.config.translate_fn:
            self._log_phase('phase_1_translate', 0, {"status": "no_function"})
            return {}
        texts_to_translate = self._extract_texts(resolved)
        try:
            translation_map = self.config.translate_fn(texts_to_translate)
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_1_translate', duration, {"status": "success", "strings_translated": len(translation_map)})
            return translation_map
        except Exception as e:
            raise RuntimeError(f"Phase 1 (Translation) failed: {e}")
        
    def _phase_2_transform_masters_layouts(self, prs, resolved):
        start_time = time.monotonic()
        try:
            registry = TemplateRegistry()
            transformer = MasterLayoutTransformer(prs, registry)
            transformer.transform_all_masters()
            report = transformer.transform_all_layouts()
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_2_transform_masters', duration, {"status": "success"})
            return report
        except NameError as e:
            logger.warning(f"Phase 2 stubbed: {e}")
            self._log_phase('phase_2_transform_masters', 0, {"status": "stubbed"})
            return None
        
    def _phase_3_transform_slides(self, prs, resolved, translations):
        start_time = time.monotonic()
        try:
            analyzer = LayoutAnalyzer()
            layout_classifications = analyzer.classify_slides(prs)
            transformer = SlideContentTransformer(prs, layout_classifications, translations)
            report = transformer.transform_all_slides()
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_3_transform_slides', duration, {"status": "success"})
            return report
        except NameError as e:
            logger.warning(f"Phase 3 stubbed: {e}")
            self._log_phase('phase_3_transform_slides', 0, {"status": "stubbed"})
            return None
        
    def _phase_4_typography(self, prs):
        start_time = time.monotonic()
        try:
            normalizer = TypographyNormalizer(prs, max_reduction_pct=self.config.max_font_reduction_pct)
            report = normalizer.normalize_all()
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_4_typography', duration, {"status": "success"})
            return report
        except NameError as e:
            logger.warning(f"Phase 4 stubbed: {e}")
            self._log_phase('phase_4_typography', 0, {"status": "stubbed"})
            return None
        
    def _phase_5_validate(self, prs, resolved):
        start_time = time.monotonic()
        try:
            validator = StructuralValidator(prs, resolved)
            report = validator.validate()
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_5_validate', duration, {"status": "success", "passed": report.passed})
            return report
        except NameError as e:
            logger.warning(f"Phase 5 stubbed: {e}")
            self._log_phase('phase_5_validate', 0, {"status": "stubbed"})
            return None
        
    def _extract_texts(self, resolved):
        texts = []
        if not resolved:
            return texts
        try:
            for slide in resolved.slides:
                for shape in slide.shapes:
                    for para in shape.paragraphs:
                        para_text = "".join(run.text for run in para.runs if run.text).strip()
                        if para_text:
                            texts.append(para_text)
            seen = set()
            unique_texts = []
            for t in texts:
                if t not in seen:
                    seen.add(t)
                    unique_texts.append(t)
            return unique_texts
        except AttributeError:
            return []
        
    def _log_phase(self, phase_name, duration_ms, report):
        self._phase_reports[phase_name] = {"duration_ms": duration_ms, "report": report}
        logger.debug(f"{phase_name} completed in {duration_ms:.1f}ms: {report}")
