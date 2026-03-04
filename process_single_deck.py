#!/usr/bin/env python3
"""
Process a single PPTX deck through the SlideShift v2 pipeline with translations.

Usage:
  # With pre-generated cache:
  python process_single_deck.py <input.pptx> <output.pptx> <translations.json>

  # With LLM translation (GPT-5.2 + Claude Sonnet 4.6 QA):
  python process_single_deck.py <input.pptx> <output.pptx> --llm-translate

Flags:
  --llm-translate     Use dual-LLM translation (requires OPENAI_API_KEY env var,
                      optionally ANTHROPIC_API_KEY for QA pass)
  --no-qa             Skip Claude QA pass (only with --llm-translate)
  --no-vqa            Skip Phase 6 visual QA
  --issue-log PATH    Path for JSON Lines issue log

IMPORTANT: Must monkey-patch python-pptx for ZIP_STORED to handle large files.
"""

import json
import logging
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

# Add parent dir to path
sys.path.insert(0, str(Path(__file__).resolve().parent))
sys.path.insert(0, str(Path(__file__).resolve().parent / "skills" / "pptx" / "scripts"))

logging.basicConfig(level=logging.WARNING, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)

# ── Monkey-patch python-pptx for ZIP_STORED saves ───────────────────────────
try:
    import pptx.opc.phys_pkg as _phys_pkg  # older versions
except ModuleNotFoundError:
    _phys_pkg = None
import pptx.opc.serialized as _ser

_orig_zipf = None

class _ZipStoredPatch:
    """Patch ZipPkgWriter to use ZIP_STORED compression."""
    @staticmethod
    def apply():
        from pptx.util import lazyproperty
        def _zipf_stored(self):
            return zipfile.ZipFile(self._pkg_file, 'w', zipfile.ZIP_STORED)
        _ser._ZipPkgWriter._zipf = lazyproperty(_zipf_stored)
        print("  [Patch] ZIP_STORED monkey-patch applied")

_ZipStoredPatch.apply()

from pptx import Presentation

from slideshift_v2.rtl_transforms import MasterLayoutTransformer, SlideContentTransformer
from slideshift_v2.layout_analyzer import LayoutAnalyzer
from slideshift_v2.template_registry import TemplateRegistry
from slideshift_v2.typography import TypographyNormalizer
from slideshift_v2.property_resolver import PropertyResolver
from slideshift_v2.visual_qa import run_vqa
from slideshift_v2.smartart_translator import translate_smartart_in_pptx


def recompress_pptx(pptx_path: Path) -> None:
    """Re-compress a ZIP_STORED pptx to ZIP_DEFLATED to reduce file size."""
    import tempfile
    tmp = pptx_path.with_suffix('.pptx.tmp')
    try:
        with zipfile.ZipFile(str(pptx_path), 'r') as zin:
            with zipfile.ZipFile(str(tmp), 'w', zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    zout.writestr(item, data)
        shutil.move(str(tmp), str(pptx_path))
        print(f"  [Recompress] Saved compressed: {pptx_path.stat().st_size / 1024:.0f} KB")
    except Exception as e:
        print(f"  [WARN] Recompress failed: {e}")
        if tmp.exists():
            tmp.unlink()


def _extract_all_texts(prs) -> list:
    """Extract all text strings from a presentation for LLM translation."""
    from pptx.util import Inches
    texts = set()
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and len(text) > 0:
                        texts.add(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text and len(text) > 0:
                            texts.add(text)
    
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.add(text)
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.add(text)
    
    return sorted(texts)


def process_deck(input_path: str, output_path: str, translations_json: str = None,
                 enable_vqa: bool = True, issue_log_path: str = None,
                 llm_translate: bool = False, enable_qa: bool = True):
    """Process a single deck with full pipeline + translations."""
    input_p = Path(input_path)
    output_p = Path(output_path)
    
    print(f"\n{'='*60}")
    print(f"  Processing: {input_p.name}")
    print(f"{'='*60}")
    
    translations = {}
    
    if llm_translate:
        import os
        from slideshift_v2.llm_translator import DualLLMTranslator, TranslatorConfig
        
        openai_key = os.environ.get('OPENAI_API_KEY', '')
        anthropic_key = os.environ.get('ANTHROPIC_API_KEY', '')
        
        if not openai_key:
            print("  [ERROR] OPENAI_API_KEY environment variable required for --llm-translate")
            return False
        
        if not anthropic_key:
            print("  [WARN] ANTHROPIC_API_KEY not set \u2014 Claude QA pass will be skipped")
        
        cache_dir = output_p.parent / 'translations_cache_llm'
        cache_dir.mkdir(parents=True, exist_ok=True)
        cache_path = str(cache_dir / f"{input_p.stem}_llm.json")
        
        print("  Phase 1a: Extracting texts from presentation...")
        temp_prs = Presentation(str(input_p))
        all_texts = _extract_all_texts(temp_prs)
        print(f"  Phase 1a: Extracted {len(all_texts)} unique text strings")
        del temp_prs
        
        print(f"  Phase 1b: Translating via GPT + Claude QA...")
        config = TranslatorConfig(
            openai_api_key=openai_key,
            anthropic_api_key=anthropic_key,
            enable_qa_pass=enable_qa and bool(anthropic_key),
        )
        translator = DualLLMTranslator(config)
        translations = translator.translate(all_texts, cache_path=cache_path)
        
        report = translator.report
        print(f"  Phase 1b: {report.translated} translated, "
              f"{report.from_cache} from cache, "
              f"{report.qa_issues_found} QA issues found, "
              f"{report.qa_issues_fixed} fixed")
        if report.errors:
            for err in report.errors[:3]:
                print(f"    ERROR: {err}")
        print(f"  Phase 1b: Estimated cost: ${report.estimated_cost_usd():.4f}")
        print(f"  Phase 1b: Elapsed: {report.elapsed_seconds:.1f}s")
        
    elif translations_json:
        trans_p = Path(translations_json)
        if trans_p.exists():
            with open(trans_p, 'r', encoding='utf-8') as f:
                translations = json.load(f)
            print(f"  Loaded {len(translations)} translations from {trans_p.name}")
        else:
            print(f"  [WARN] No translations file: {trans_p}")
    else:
        print(f"  [WARN] No translation source specified (use --llm-translate or provide translations.json)")
    
    shutil.copy2(str(input_p), str(output_p))
    
    prs = Presentation(str(output_p))
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    print(f"  Slide dimensions: {slide_width}x{slide_height} EMU, {len(prs.slides)} slides")
    
    layout_classifications = {}
    try:
        analyzer = LayoutAnalyzer(prs)
        layout_classifications = analyzer.analyze_all()
        print(f"  Phase 0: Analyzed {len(layout_classifications)} slide layouts")
    except Exception as e:
        print(f"  [WARN] Layout analysis: {e}")
    
    registry = None
    try:
        registry = TemplateRegistry(slide_width, slide_height)
    except Exception as e:
        print(f"  [WARN] Template registry: {e}")
    
    try:
        ml_transformer = MasterLayoutTransformer(prs, registry)
        master_report = ml_transformer.transform_all_masters()
        layout_report = ml_transformer.transform_all_layouts()
        print(f"  Phase 2: {master_report.total_changes} master + {layout_report.total_changes} layout changes")
    except Exception as e:
        print(f"  [ERROR] Phase 2: {e}")
    
    try:
        content_transformer = SlideContentTransformer(
            prs,
            template_registry=registry,
            layout_classifications=layout_classifications,
            translations=translations,
        )
        slide_report = content_transformer.transform_all_slides()
        print(f"  Phase 3: {slide_report.total_changes} slide changes")
        if slide_report.errors:
            for err in slide_report.errors[:5]:
                print(f"    ERROR: {err}")
    except Exception as e:
        print(f"  [ERROR] Phase 3: {e}")
        import traceback; traceback.print_exc()
    
    try:
        normalizer = TypographyNormalizer(prs)
        typo_report = normalizer.normalize_all()
        print(f"  Phase 4: {typo_report.total_changes} typography changes")
    except Exception as e:
        print(f"  [WARN] Phase 4: {e}")
    
    try:
        prs.save(str(output_p))
        size_kb = output_p.stat().st_size / 1024
        print(f"  Saved: {output_p.name} ({size_kb:.0f} KB)")
    except Exception as e:
        print(f"  [ERROR] Save failed: {e}")
        return False
    
    try:
        recompress_pptx(output_p)
    except Exception as e:
        print(f"  [WARN] Recompress: {e}")
    
    vqa_report = None
    if enable_vqa:
        try:
            log_path = issue_log_path or str(
                output_p.parent / "vqa_issues.jsonl"
            )
            print(f"  Phase 6: Running closed-loop VQA...")
            vqa_report = run_vqa(
                original_pptx=str(input_p),
                converted_pptx=str(output_p),
                max_slides=20,
                issue_log_path=log_path,
                deck_name=input_p.name,
            )
            print(f"  Phase 6: {vqa_report.summary()}")
            if vqa_report.remediation_attempted > 0:
                print(f"    Remediated: {vqa_report.remediation_attempted} slides, "
                      f"{vqa_report.remediation_successful} improved")
            if vqa_report.issues_logged > 0:
                print(f"    Issues logged: {vqa_report.issues_logged} \u2192 {log_path}")
        except Exception as e:
            print(f"  [WARN] Phase 6 VQA: {e}")

    print(f"  DONE: {input_p.name}")
    return True


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <output.pptx> [translations.json]")
        print(f"       {sys.argv[0]} <input.pptx> <output.pptx> --llm-translate [--no-qa]")
        print(f"")
        print(f"  --llm-translate  Use GPT + Claude for translation (requires OPENAI_API_KEY)")
        print(f"  --no-qa          Skip Claude QA pass (only with --llm-translate)")
        print(f"  --no-vqa         Skip Phase 6 visual QA")
        print(f"  --issue-log PATH Path for JSON Lines issue log")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    llm_translate = '--llm-translate' in sys.argv
    enable_qa = '--no-qa' not in sys.argv
    enable_vqa = '--no-vqa' not in sys.argv
    
    translations_json = None
    if len(sys.argv) > 3 and not sys.argv[3].startswith('--'):
        translations_json = sys.argv[3]
    
    issue_log = None
    for i, arg in enumerate(sys.argv):
        if arg == '--issue-log' and i + 1 < len(sys.argv):
            issue_log = sys.argv[i + 1]
    
    success = process_deck(
        input_file, output_file,
        translations_json=translations_json,
        enable_vqa=enable_vqa,
        issue_log_path=issue_log,
        llm_translate=llm_translate,
        enable_qa=enable_qa,
    )
    sys.exit(0 if success else 1)
